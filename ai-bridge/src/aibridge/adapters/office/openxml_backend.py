"""
OpenXML Backend — Phase III v0.11.0

Pure Python cross-platform Office backend using:
- python-docx for Word (.docx)
- openpyxl for Excel (.xlsx)
- python-pptx for PowerPoint (.pptx)

No COM, no LibreOffice subprocess — works on Windows/Mac/Linux.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

import pandas as pd

from aibridge.adapters.office.base import (
    DocumentHandle,
    OfficeAdapter,
    WordAdapter,
    ExcelAdapter,
    PPTAdapter,
)

logger = logging.getLogger(__name__)

BACKEND_TYPE = "openxml"


# ============ Word ============

class OpenXMLWordAdapter(WordAdapter):
    """Word adapter using python-docx (pure Python, cross-platform)."""

    _backend_type = BACKEND_TYPE

    @classmethod
    def is_available(cls) -> bool:
        try:
            import docx  # noqa: F401
            return True
        except ImportError:
            return False

    def open_document(self, path: Path) -> DocumentHandle:
        import docx

        path = Path(path)
        if path.exists():
            self._doc = docx.Document(str(path))
        else:
            self._doc = docx.Document()
        handle = DocumentHandle(path=path, backend_type=BACKEND_TYPE)
        self._open_handles[str(path)] = handle
        return handle

    def save_as(
        self, doc: DocumentHandle, target_path: Path, format: str | None = None
    ) -> DocumentHandle:
        target_path = Path(target_path)
        if format:
            target_path = target_path.with_suffix(f".{format}")
        self._doc.save(str(target_path))
        return DocumentHandle(path=target_path, backend_type=BACKEND_TYPE)

    def export_pdf(
        self, doc: DocumentHandle, output_path: Path | None = None
    ) -> Path:
        # OpenXML backend cannot natively export PDF.
        # Save as docx first; PDF conversion requires LibreOffice or win32.
        output_path = Path(output_path or doc.path.with_suffix(".pdf"))
        logger.warning(
            "OpenXML backend cannot natively export PDF. "
            "Saving as docx. Use LibreOffice or Win32 backend for PDF."
        )
        self._doc.save(str(output_path.with_suffix(".docx")))
        return output_path

    def extract_text(self, doc: DocumentHandle) -> str:
        paragraphs = [p.text for p in self._doc.paragraphs]
        return "\n".join(paragraphs)

    def extract_tables(self, doc: DocumentHandle) -> list[pd.DataFrame]:
        tables = []
        for table in self._doc.tables:
            data = []
            for row in table.rows:
                data.append([cell.text for cell in row.cells])
            if data:
                tables.append(pd.DataFrame(data[1:], columns=data[0])
                              if len(data) > 1 and data[0]
                              else pd.DataFrame(data))
        return tables

    def close(self, doc: DocumentHandle) -> None:
        self._doc = None
        self._open_handles.pop(str(doc.path), None)

    def insert_text(
        self, doc: DocumentHandle, text: str, position: str = "end"
    ) -> None:
        if position == "start":
            if self._doc.paragraphs:
                self._doc.paragraphs[0].insert_paragraph_before(text)
            else:
                self._doc.add_paragraph(text)
        else:
            self._doc.add_paragraph(text)

    def replace_text(self, doc: DocumentHandle, old: str, new: str) -> int:
        count = 0
        for para in self._doc.paragraphs:
            if old in para.text:
                # python-docx: inline runs make replacement tricky.
                # Simple approach: replace in paragraph text.
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        count += 1
        return count

    def add_table(
        self,
        doc: DocumentHandle,
        rows: int,
        cols: int,
        data: list[list] | None = None,
    ) -> int:
        table = self._doc.add_table(rows=rows, cols=cols)
        if data:
            for i, row_data in enumerate(data):
                for j, cell_val in enumerate(row_data):
                    if i < rows and j < cols:
                        table.cell(i, j).text = str(cell_val)
        return len(self._doc.tables) - 1

    def get_comments(self, doc: DocumentHandle) -> list[dict]:
        # python-docx does not support comments extraction natively
        return []


# ============ Excel ============

class OpenXMLExcelAdapter(ExcelAdapter):
    """Excel adapter using openpyxl (pure Python, cross-platform)."""

    _backend_type = BACKEND_TYPE

    @classmethod
    def is_available(cls) -> bool:
        try:
            import openpyxl  # noqa: F401
            return True
        except ImportError:
            return False

    def open_document(self, path: Path) -> DocumentHandle:
        import openpyxl

        path = Path(path)
        if path.exists():
            self._wb = openpyxl.load_workbook(str(path))
        else:
            self._wb = openpyxl.Workbook()
        handle = DocumentHandle(path=path, backend_type=BACKEND_TYPE)
        self._open_handles[str(path)] = handle
        return handle

    def save_as(
        self, doc: DocumentHandle, target_path: Path, format: str | None = None
    ) -> DocumentHandle:
        target_path = Path(target_path)
        if format:
            target_path = target_path.with_suffix(f".{format}")
        self._wb.save(str(target_path))
        return DocumentHandle(path=target_path, backend_type=BACKEND_TYPE)

    def export_pdf(
        self, doc: DocumentHandle, output_path: Path | None = None
    ) -> Path:
        output_path = Path(output_path or doc.path.with_suffix(".pdf"))
        logger.warning(
            "OpenXML backend cannot natively export PDF. "
            "Saving as xlsx. Use LibreOffice or Win32 backend for PDF."
        )
        self._wb.save(str(output_path.with_suffix(".xlsx")))
        return output_path

    def extract_text(self, doc: DocumentHandle) -> str:
        texts = []
        for ws in self._wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(c) for c in row if c is not None
                )
                if row_text:
                    texts.append(row_text)
        return "\n".join(texts)

    def extract_tables(self, doc: DocumentHandle) -> list[pd.DataFrame]:
        tables = []
        for ws in self._wb.worksheets:
            data = list(ws.iter_rows(values_only=True))
            if data:
                df = pd.DataFrame(data[1:], columns=data[0]) if len(data) > 1 else pd.DataFrame(data)
                tables.append(df)
        return tables

    def close(self, doc: DocumentHandle) -> None:
        self._wb = None
        self._open_handles.pop(str(doc.path), None)

    def read_range(self, doc: DocumentHandle, range_str: str) -> pd.DataFrame:
        import openpyxl.utils

        sheet_name, range_part = "Sheet", range_str
        if "!" in range_str:
            sheet_name, range_part = range_str.split("!", 1)

        ws = self._wb[sheet_name] if sheet_name in self._wb.sheetnames else self._wb.active
        try:
            min_col, min_row, max_col, max_row = openpyxl.utils.range_boundaries(range_part)
            data = []
            for row in ws.iter_rows(
                min_row=min_row, max_row=max_row,
                min_col=min_col, max_col=max_col,
                values_only=True,
            ):
                data.append(list(row))
            return pd.DataFrame(data) if data else pd.DataFrame()
        except Exception:
            # Fallback: read entire sheet
            return pd.DataFrame(ws.values)

    def write_range(
        self, doc: DocumentHandle, range_str: str, data: pd.DataFrame
    ) -> None:
        import openpyxl.utils

        sheet_name, range_part = "Sheet", range_str
        if "!" in range_str:
            sheet_name, range_part = range_str.split("!", 1)

        ws = self._wb[sheet_name] if sheet_name in self._wb.sheetnames else self._wb.active
        min_col, min_row, _, _ = openpyxl.utils.range_boundaries(range_part)

        for i, row in enumerate(data.itertuples(index=False)):
            for j, val in enumerate(row):
                ws.cell(row=min_row + i, column=min_col + j, value=val)

    def add_chart(
        self,
        doc: DocumentHandle,
        data_range: str,
        chart_type: str,
        position: str,
    ) -> None:
        from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference

        chart_classes = {
            "bar": BarChart,
            "line": LineChart,
            "pie": PieChart,
            "scatter": ScatterChart,
        }
        chart_cls = chart_classes.get(chart_type, BarChart)

        ws = self._wb.active
        chart = chart_cls()
        ref = Reference(ws, range_string=data_range)
        chart.add_data(ref)
        ws.add_chart(chart, position)

    def create_pivot(
        self,
        doc: DocumentHandle,
        source_range: str,
        rows: list[str],
        values: list[str],
    ) -> None:
        # openpyxl pivot tables require significant setup.
        # For cross-platform baseline, raise NotImplementedError.
        raise NotImplementedError(
            "Pivot tables not supported in OpenXML backend. "
            "Use Win32 backend on Windows."
        )

    def apply_formula(
        self, doc: DocumentHandle, cell: str, formula: str
    ) -> None:
        import openpyxl.utils

        ws = self._wb.active
        col_letter, row_num = "", ""
        for c in cell:
            if c.isalpha():
                col_letter += c
            else:
                row_num += c
        col_idx = openpyxl.utils.column_index_from_string(col_letter)
        ws.cell(row=int(row_num), column=col_idx, value=formula)


# ============ PowerPoint ============

class OpenXMLPPTAdapter(PPTAdapter):
    """PowerPoint adapter using python-pptx (pure Python, cross-platform)."""

    _backend_type = BACKEND_TYPE

    @classmethod
    def is_available(cls) -> bool:
        try:
            import pptx  # noqa: F401
            return True
        except ImportError:
            return False

    def open_document(self, path: Path) -> DocumentHandle:
        from pptx import Presentation

        path = Path(path)
        if path.exists():
            self._prs = Presentation(str(path))
        else:
            self._prs = Presentation()
        handle = DocumentHandle(path=path, backend_type=BACKEND_TYPE)
        self._open_handles[str(path)] = handle
        return handle

    def save_as(
        self, doc: DocumentHandle, target_path: Path, format: str | None = None
    ) -> DocumentHandle:
        target_path = Path(target_path)
        if format:
            target_path = target_path.with_suffix(f".{format}")
        self._prs.save(str(target_path))
        return DocumentHandle(path=target_path, backend_type=BACKEND_TYPE)

    def export_pdf(
        self, doc: DocumentHandle, output_path: Path | None = None
    ) -> Path:
        output_path = Path(output_path or doc.path.with_suffix(".pdf"))
        logger.warning(
            "OpenXML backend cannot natively export PDF. "
            "Saving as pptx. Use LibreOffice or Win32 backend for PDF."
        )
        self._prs.save(str(output_path.with_suffix(".pptx")))
        return output_path

    def extract_text(self, doc: DocumentHandle) -> str:
        texts = []
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return "\n".join(texts)

    def extract_tables(self, doc: DocumentHandle) -> list[pd.DataFrame]:
        tables = []
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    data = []
                    for row in table.rows:
                        data.append([cell.text for cell in row.cells])
                    if data:
                        tables.append(
                            pd.DataFrame(data[1:], columns=data[0])
                            if len(data) > 1 and data[0]
                            else pd.DataFrame(data)
                        )
        return tables

    def close(self, doc: DocumentHandle) -> None:
        self._prs = None
        self._open_handles.pop(str(doc.path), None)

    def add_slide(self, doc: DocumentHandle, layout: str = "blank") -> int:
        from pptx.util import Inches

        layout_map = {
            "blank": 6,     # BLANK layout index
            "title": 0,     # TITLE layout index
            "content": 1,   # TITLE_AND_CONTENT
        }
        layout_idx = layout_map.get(layout, 6)
        slide_layout = self._prs.slide_layouts[layout_idx]
        self._prs.slides.add_slide(slide_layout)
        return len(self._prs.slides)

    def add_text_box(
        self,
        doc: DocumentHandle,
        slide_index: int,
        text: str,
        position: tuple[float, float, float, float] | None = None,
    ) -> None:
        from pptx.util import Inches

        if slide_index < 1 or slide_index > len(self._prs.slides):
            raise IndexError(f"Slide {slide_index} out of range")

        slide = self._prs.slides[slide_index - 1]
        left, top, width, height = position or (1.0, 1.0, 8.0, 2.0)
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        txBox.text_frame.text = text

    def add_image(
        self,
        doc: DocumentHandle,
        slide_index: int,
        image_path: Path,
        position: tuple[float, float, float, float] | None = None,
    ) -> None:
        from pptx.util import Inches

        if slide_index < 1 or slide_index > len(self._prs.slides):
            raise IndexError(f"Slide {slide_index} out of range")

        slide = self._prs.slides[slide_index - 1]
        left, top, width, height = position or (1.0, 1.0, 4.0, 3.0)
        slide.shapes.add_picture(
            str(image_path),
            Inches(left), Inches(top),
            Inches(width), Inches(height),
        )
